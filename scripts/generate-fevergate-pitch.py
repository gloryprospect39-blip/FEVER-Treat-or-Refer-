"""Generate FeverGate 5-minute pitch deck (.pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDES = [
    {
        "title": "FeverGate",
        "subtitle": "Treat or refer — in under a minute, without a lab",
        "bullets": [
            "Next.js PWA · Myanmar UI · deterministic clinical rules",
            "Built for community health workers at border clinics",
            "Live: fever-treat-or-refer.vercel.app",
        ],
        "notes": (
            "Good morning/afternoon. I'm pitching FeverGate — a bedside decision aid "
            "for febrile patients. No lab, no specialist on call, one clear action in "
            "under sixty seconds."
        ),
    },
    {
        "title": "The problem",
        "subtitle": "Every fever is a high-stakes fork",
        "bullets": [
            "Health worker alone · spotty connectivity · unpredictable drug stock",
            "IMCI protocols on paper ≠ usable under bedside pressure",
            "Treat severe illness at home → lives lost",
            "Refer every mild fever → referral capacity burned",
            "Adults with comorbidities need logic beyond pediatric IMCI alone",
        ],
        "notes": (
            "A worker has roughly two minutes with a febrile patient — child or adult — "
            "to decide treat here or refer. Wrong either way costs lives or scarce "
            "transport slots."
        ),
    },
    {
        "title": "The solution",
        "subtitle": "Four clear groups — one primary action",
        "bullets": [
            "A — Immediate refer · B — Urgent hospital refer",
            "C — Treat & monitor · D — Treat",
            "Named reason + drug plan + teleconsult: call now or book",
            "Stock- and endemicity-aware ACT / paracetamol plans",
            "Same labels on result cards, reports, and supervisor charts",
        ],
        "notes": (
            "Village, age band, fever, danger signs, optional vitals — assess — one "
            "full-screen card in Group A, B, C, or D with exactly what to do next."
        ),
    },
    {
        "title": "How it works",
        "subtitle": "Under 60 seconds per patient",
        "bullets": [
            "1. Clinic context — malaria endemicity + start-of-day stock check-in",
            "2. Patient — village (A–G), age band, danger signs, optional vitals",
            "3. Assess → Group A–D card → act → encounter log → next patient",
            "Age 15 splits pediatric (lean IMCI) vs adult (comorbidities)",
            "Medics stay on triage; supervisor hub is URL-only (no clutter)",
        ],
        "notes": (
            "Morning setup once per shift — endemicity and stock baseline. Then "
            "repeat: assess, act, log, next. Supervisors open /supervisor separately."
        ),
    },
    {
        "title": "Safety first",
        "subtitle": "Deterministic rules — no AI in the decision path",
        "bullets": [
            "Any positive IMCI danger sign → refer pathway (any age)",
            "Neonate <2 months + fever → immediate referral (Group A)",
            "Adults: qSOFA & NEWS2 when vitals available",
            "Unified danger-sign and comorbidity icons — less scan friction",
            "Success bar: zero false negatives on danger signs",
        ],
        "notes": (
            "This is a safety product. Pure TypeScript rules — no LLM decides. "
            "Auditable and testable."
        ),
    },
    {
        "title": "What we shipped",
        "subtitle": "Field-ready — including this week's UX",
        "bullets": [
            "Myanmar PWA · villages A–G · visit trace · printable referral form",
            "Start-of-day clinic stock check-in + per-patient dispensing log",
            "Supervisor hub: KPIs, decision mix, village volume, drug & activity",
            "Groups A–D wording · teleconsult CTAs · recommendation-style forms",
            "Shared decision labels across triage, reports, and dashboard",
        ],
        "notes": (
            "Built for real clinic friction — stock unpredictability, epidemiology "
            "reporting, and a separate supervisor view so medics stay focused."
        ),
    },
    {
        "title": "Live demo beats",
        "subtitle": "~30 seconds each",
        "bullets": [
            "A/B: under-five convulsions → refer + call teleconsultation now",
            "D: uncomplicated fever, ACT in stock → presumptive ACT plan",
            "C: treat & monitor → follow-up return wording without referral noise",
            "Supervisor: decision mix + village volume + ACT stock-outs",
        ],
        "notes": (
            "If demoing live, run these scenarios. Close with: "
            "'It told me exactly what to do — and why.'"
        ),
    },
    {
        "title": "Why now & ask",
        "subtitle": "From prototype to pilot",
        "bullets": [
            "v1: Groups A–D + stock check-in + supervisor reporting",
            "Next: MOH clinical sign-off · teleconsult routing · district sync",
            "Ask: field validation partners + epidemiology feedback",
            "",
            '"I picked the age band, tapped the signs, and it told me exactly '
            'what to do — Group A through D, with exactly why."',
        ],
        "notes": (
            "FeverGate turns guideline knowledge into bedside action when the worker is "
            "alone and stock is unpredictable. Thank you — questions welcome."
        ),
    },
]

TEAL = RGBColor(0x0D, 0x94, 0x88)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROSE = RGBColor(0xBE, 0x12, 0x3C)


def add_title_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "FeverGate — 5-Minute Pitch"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(8.4), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Border clinic febrile triage · Treat or refer"
    sp.font.size = Pt(18)
    sp.font.color.rgb = MUTED
    sp.alignment = PP_ALIGN.CENTER

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.5))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "Updated · Groups A–D · stock check-in · supervisor hub"
    tp.font.size = Pt(14)
    tp.font.color.rgb = TEAL
    tp.alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.4))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "9 slides · ~5 minutes total"
    fp.font.size = Pt(12)
    fp.font.color.rgb = MUTED
    fp.alignment = PP_ALIGN.CENTER


def add_content_slide(prs: Presentation, data: dict, index: int) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # slide number
    num = slide.shapes.add_textbox(Inches(9.0), Inches(0.35), Inches(0.8), Inches(0.3))
    np = num.text_frame.paragraphs[0]
    np.text = str(index)
    np.font.size = Pt(11)
    np.font.color.rgb = MUTED
    np.alignment = PP_ALIGN.RIGHT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(8.2), Inches(0.9))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = data["title"]
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = SLATE

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(8.2), Inches(0.5))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = data["subtitle"]
    sp.font.size = Pt(16)
    sp.font.color.rgb = TEAL

    body = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(8.0), Inches(4.2))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for bullet in data["bullets"]:
        if not bullet:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        is_quote = bullet.startswith('"')
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18 if is_quote else 16)
        p.font.color.rgb = ROSE if is_quote else SLATE
        p.font.italic = is_quote
        p.space_after = Pt(8)

    notes = slide.notes_slide.notes_text_frame
    notes.text = data["notes"]


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    for i, slide_data in enumerate(SLIDES, start=1):
        add_content_slide(prs, slide_data, i)

    out = Path(__file__).resolve().parents[1] / "FeverGate-5min-Pitch.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
